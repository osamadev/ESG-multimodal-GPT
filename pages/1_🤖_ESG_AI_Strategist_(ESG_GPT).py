import vertexai
from vertexai.preview.generative_models import GenerativeModel, Part
import streamlit as st
import base64
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.llms import OpenAI
from langchain.vectorstores import Pinecone
import re
import os
import pinecone
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import textwrap
import io
from langchain.agents import Tool, initialize_agent
from langchain.chains import LLMChain, RetrievalQA
from langchain.memory import ConversationBufferWindowMemory
from langchain.llms import VertexAI
from langchain.prompts import PromptTemplate
from openai import OpenAI
from streamlit_extras.switch_page_button import switch_page

from google.oauth2 import service_account

def init_app():
    os.environ["GOOGLE_API_KEY"] = st.secrets["GOOGLE_API_KEY"]
    os.environ["OPENAI_API_KEY"] = st.secrets["OPENAI_API_KEY"]
    os.environ["PINECONE_API_KEY"] = st.secrets["PINECONE_API_KEY"]
    os.environ["PINECONE_ENV"] = st.secrets["PINECONE_ENV"]

    project_id = "vital-future-408219"
    location = "us-central1"
    credentials = service_account.Credentials.from_service_account_info(st.secrets.connections_gcs)
    vertexai.init(project=project_id, location=location, credentials=credentials)

class PresentationCreationTool(Tool):
    def __init__(self, name, description):
        super().__init__(name=name, description=description, func=self.create_presentation)

    def create_presentation(self, conversation_history):
        # Generate slide content with LLM
        from openai import OpenAI
        OpenAI.api_key = st.secrets["OPENAI_API_KEY"]
        client = OpenAI()

        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a PowerPoint expert, designer, content organizer, and slides generator."},
                {"role": "user", "content": f"""Format and generate PowerPoint slides using this context:\n{conversation_history}\n
                 Prepare the slides from the given conversation history and break the content logically according to the topic into slides, 
                 each slide should start with a title in that format "slide 1: Title of slide 1" in a separate line. 
                 The title of each slide should be in the required format and in a new line"""}
            ],
            max_tokens=2000,
            temperature=0.5
        )
        generated_content = response.choices[0].message.content

        # Create a presentation object
        prs = Presentation()

        # Define a professional and elegant design
        bg_color = RGBColor(245, 245, 245)  # Light grey background for elegance
        title_font_color = RGBColor(0, 51, 102)  # Dark blue for titles
        content_font_color = RGBColor(32, 32, 32)  # Almost black for content
        font_family = "Arial"  # Professional font family

        # Updated regex to correctly extract slides
        slides = re.findall(r'Slide \d+: (.*?)\n(?=Slide \d+:|$)', generated_content, re.DOTALL)

        for slide in slides:
            # Split title and content
            title, content = slide.split('\n', 1)

            # Add a blank slide
            slide_layout = prs.slide_layouts[6]  # 6 corresponds to blank slide layout
            slide = prs.slides.add_slide(slide_layout)

            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

            # Add title
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.clear()  # Clear any existing content in the frame
            title_frame.word_wrap = True
            title_p = title_frame.add_paragraph()
            title_p.text = title.strip()
            title_p.font.size = Pt(36)
            title_p.font.color.rgb = title_font_color
            title_p.font.name = font_family

            # Add content as bullet list
            content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            content_frame = content_shape.text_frame
            content_frame.clear()  # Clear any existing content in the frame
            content_frame.word_wrap = True
            for line in content.split('\n'):
                p = content_frame.add_paragraph()
                p.text = line.strip()
                p.level = 0  # Adjust as needed for nested bullets
                p.font.size = Pt(20)
                p.font.color.rgb = content_font_color
                p.font.name = font_family

            # Adjust margins
            content_frame.margin_left = Emu(457200)  # Approx. 0.5 inch
            content_frame.margin_right = Emu(457200)  # Approx. 0.5 inch
            content_frame.margin_top = Emu(228600)  # Approx. 0.25 inch
            content_frame.margin_bottom = Emu(228600)  # Approx. 0.25 inch

        # Save to a BytesIO object
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        return ppt_io, generated_content
    
def preview_presentation_content(generated_content):
    slides = re.findall(r'Slide \d+: (.*?)\n(?=Slide \d+:|$)', generated_content, re.DOTALL)
    expander_title = "**📑 Preview Slides Content**"  
    with st.sidebar.expander(expander_title, expanded=True):
        for i, slide in enumerate(slides, 1):
            title, content = slide.split('\n', 1)
            st.subheader(f"Slide {i}: {title.strip()}")
            st.write(content.strip())
            st.markdown("---")

# Initialize the tool
presentation_tool = PresentationCreationTool(
    name="Presentation Generation Tool",
    description="Creates and generates PowerPoint presentations from conversation history. It will be used when the user requests to generate a presnetation out of his conversation"
)

class ImageExplanationTool(Tool):
    def __init__(self, name, description, model_name="gemini-pro-vision"):
        # Initialize the GenerativeModel
        # self.model = GenerativeModel(model_name)

        # Define a lambda function that captures 'self' and calls 'run'
        run_func = lambda image_data, mime_type, combined_prompt: self.run(image_data, mime_type, combined_prompt)

        # Pass the lambda function as the 'func' argument to the Tool base class
        super().__init__(name=name, description=description, func=run_func)

    def run(self, image_data, mime_type, combined_prompt):
        if image_data is not None:
            image_part = Part.from_data(data=base64.b64decode(image_data), mime_type=mime_type)
            model = GenerativeModel("gemini-pro-vision")
            responses = model.generate_content(
                [image_part, combined_prompt],
                generation_config={
                    "max_output_tokens": 2048,
                    "temperature": 0.4,
                    "top_p": 1,
                    "top_k": 32
                },
                stream=True,
            )
        else:
            responses = []

        result = ""
        for response in responses:
            result += response.candidates[0].content.parts[0].text + "\n"
        return result

# Initialize your ImageExplanationTool
image_explanation_tool = ImageExplanationTool(
    name="Image Explanation Tool",
    description="Explains or describes images based on user prompts.", 
)

class DalleImageGenerationTool(Tool):
    def __init__(self, name, description):
        super().__init__(name=name, description=description, func=self.generate_image)

    def generate_image(self, prompt):
        try:
            client = OpenAI()
            response = client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )
            image_url = response.data[0].url
            return image_url
        except Exception as e:
            print(f"Error generating image: {e}")
            return None

# Initialize the tool
image_generation_tool = DalleImageGenerationTool(
    name="DALL-E Image Generation Tool",
    description="Generates images based on user prompts using DALL-E 3."
)

# Function to initialize Pinecone
def initialize_pinecone():
    index_name = 'esg-index'
    pinecone.init(api_key=os.environ["PINECONE_API_KEY"], environment=os.environ["PINECONE_ENV"])

    if index_name not in pinecone.list_indexes():
        pinecone.create_index(name=index_name, metric='cosine', dimension=1536)

    index = pinecone.Index(index_name)
    return index

def create_agent():
    index = initialize_pinecone()
    model_name = 'text-embedding-ada-002'
    openai_api_key = st.secrets["OPENAI_API_KEY"]

    embed = OpenAIEmbeddings(model=model_name, openai_api_key=openai_api_key)
    vectorstore = Pinecone(index, embed, "text")

    llm = VertexAI(model_name="gemini-pro")

    conversational_memory = ConversationBufferWindowMemory(
        memory_key='chat_history', input_key='input', output_key='output', k=5, return_messages=True)

    prompt_template = PromptTemplate(
        template="""You are an advanced ESG expert specialized in helping companies, decision makers and CEOs to 
            adopt ESG practices, policies and sustainabaility strategies. You should answer the questions 
            based on the given context:
            {context}
            Question: {question}
            Answer: """,
            input_variables=["context", "question"])
    
    qa = RetrievalQA.from_chain_type(llm=llm, chain_type="stuff", 
                                     retriever=vectorstore.as_retriever(),
                                     chain_type_kwargs={"prompt": prompt_template})

    tools = [Tool(name='ESG Knowledge Base', func=qa.run, description="""Use this tool when the user prompt has 
                  any thing related to ESG, sustainability or sustainability development goals""")]

    system_message = """
            You are an advanced ESG expert specialized in helping companies, decision makers and CEOs to 
            adopt ESG practices, policies and sustainabaility strategies. You should help them find the right answers 
            for thier ESG and sustsability-related questions to plan for ESG adoption and reporting in their companies.
        """

    prompt = PromptTemplate(
    input_variables=["query", ],
    template="{query}")

    llm_chain = LLMChain(llm=llm, prompt=prompt)

    # initialize the LLM tool
    llm_tool = Tool(
        name='Language Model',
        func=llm_chain.run,
        description='Use this tool for general purpose queries and logic')
    
    tools.append(llm_tool)

    tools.append(image_explanation_tool)

    # tools.append(image_generation_tool)

    tools.append(presentation_tool)

    return initialize_agent(agent='chat-conversational-react-description', tools=tools, llm=llm,
                            verbose=True, max_iterations=3, early_stopping_method='generate',
                            memory=conversational_memory, return_intermediate_steps=True,
                            handle_parsing_errors=True, agent_kwargs={"system_message": system_message}), qa

def generate(image_data, mime_type, combined_prompt):
    
    if image_data is not None:
        imagePart = Part.from_data(data=image_data, mime_type=mime_type)
        model = GenerativeModel("gemini-pro-vision")
        response = model.generate_content(
            [imagePart, combined_prompt],
            generation_config={
                "max_output_tokens": 2048,
                "temperature": 0.4,
                "top_p": 1,
                "top_k": 32
            },
        )
        result = response.candidates[0].content.parts[0].text
    else:
        model = GenerativeModel("gemini-pro")
        responses = model.generate_content(
            combined_prompt,
            generation_config={
                "max_output_tokens": 2048,
                "temperature": 0.4,
                "top_p": 1,
                "top_k": 32
            },
            stream=True
        )
        result = ""
        for response in responses:
            result += response.candidates[0].content.parts[0].text + "\n"

    return result

def generate_response(agent, prompt):
    try:
        response = agent({"input": prompt})
    except Exception as e:
        return f"Error has occurred with the following details: \n{e}"
    return response["output"]

if __name__ == "__main__":
    if "authentication_status" not in st.session_state \
    or st.session_state["authentication_status"] == None or st.session_state["authentication_status"] == False:
        switch_page("Home")

    init_app()
    agent, qa_chain = create_agent()

    # Streamlit UI
    st.set_page_config(page_title="🌱 ESG AI Strategist", page_icon="🌍", layout="wide")

    # Title with emojis
    st.title("🌱 ESG AI Strategist 🤖")
    st.caption("🌿🌍 ESG AI Strategist: Revolutionizing sustainable futures with AI-powered insights, guiding companies and decision-makers to embrace and excel in ESG practices and Sustainable Development Goals (SDGs). 💡🚀 Propel your business towards a greener, more responsible tomorrow!")

    with st.sidebar:
        st.markdown("## ESG Multimodal GPT")
        st.write("Navigate the path to sustainability with ESG AI Strategist App 🚀💼,where advanced GPT technology meets eco-conscious business strategies 🌱📊.")

        # Button to generate PowerPoint slides
        if st.button('Generate Slides'):
            if 'history' in st.session_state and st.session_state['history']:
                conversation_history = " ".join([text for _, text in st.session_state['history']])
                ppt_file, generated_content = presentation_tool.run(conversation_history)
                st.sidebar.download_button(
                    label="Download Slides",
                    data=ppt_file,
                    file_name="conversation_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                preview_presentation_content(generated_content)
            else:
                st.sidebar.info("No conversation history to generate presentation.")

    # Initialize session state for conversation history and file uploader state
    if 'history' not in st.session_state:
        st.session_state['history'] = []
    if 'show_file_uploader' not in st.session_state:
        st.session_state['show_file_uploader'] = False

    # Text input for the prompt
    text_prompt = st.chat_input("Enter your text prompt here:")

    # Define a callback function to update the session state
    def toggle_file_uploader():
        st.session_state['show_file_uploader'] = not st.session_state.get('show_file_uploader', False)

    # Checkbox to show/hide the file uploader
    # Using the callback to update the session state
    st.checkbox("Analyze and get insights about diagrams, infographics or charts?",
                value=st.session_state.get('show_file_uploader', False),
                on_change=toggle_file_uploader)

    # Conditionally display the file uploader based on the checkbox
    uploaded_file = None
    if st.session_state.get('show_file_uploader', False):
        uploaded_file = st.file_uploader("Upload an Image to get insights", type=["png", "jpg", "jpeg"])
        if uploaded_file is not None:
            st.image(uploaded_file, caption='', use_column_width=False, width=400)


    # Process and generate response
    if text_prompt:
        # Get MIME type of the uploaded file and base64 encode the image
        if uploaded_file is not None:
            mime_type = uploaded_file.type
            image_data = uploaded_file.getvalue()

        # Combine text prompts (text + speech-to-text)
        combined_prompt = text_prompt # Add your speech-to-text processing here

        if uploaded_file is not None:
            # Generate response
            response = generate(image_data, mime_type, combined_prompt)
        else:
            response = generate_response(agent, combined_prompt)

        # Update conversation history (prepend to display recent conversations at the top)
        st.session_state['history'].insert(0, (combined_prompt, response))

    # Display conversation history in a scrollable, stylish panel
    with st.expander("**📜 Conversation History**", expanded=True):
        # Reverse the conversation history to show the last AI response at the top
        for user_text, model_response in reversed(st.session_state['history']):
            # Styling for human user
            col1, col2 = st.columns([1, 5])
            with col1:
                st.markdown("👤 **You:**")
            with col2:
                st.info(user_text)

            # Styling for AI response
            col3, col4 = st.columns([1, 5])
            with col3:
                st.markdown("🤖 **AI:**")
            with col4:
                st.success(model_response)

            st.markdown("---")